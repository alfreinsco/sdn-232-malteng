from pathlib import Path
import sys

sys.path.insert(0, "/tmp/sisd232-ppt-deps")

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SCREENSHOTS = ROOT / "docs" / "screenshots"
OUTPUT = ROOT / "docs" / "presentasi-aplikasi-sisd232.pptx"
LOGO = ROOT / "public" / "logo-malteng.png"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
SKY = RGBColor(2, 132, 199)
SKY_DARK = RGBColor(3, 105, 161)
SKY_LIGHT = RGBColor(224, 242, 254)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)
LINE = RGBColor(226, 232, 240)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(225, 29, 72)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_background(slide, color=OFF_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color=NAVY, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=17, gap=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, (lead, detail) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.level = 0
        first = paragraph.add_run()
        first.text = lead
        first.font.name = "Aptos"
        first.font.size = Pt(size)
        first.font.bold = True
        first.font.color.rgb = NAVY
        second = paragraph.add_run()
        second.text = detail
        second.font.name = "Aptos"
        second.font.size = Pt(size)
        second.font.color.rgb = SLATE
    return box


def add_card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_badge(slide, text, x, y, w, color=SKY_DARK, fill=SKY_LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color
    return shape


def add_picture_contain(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    width = image_w * scale
    height = image_h * scale
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_picture_cover(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    frame_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if image_ratio > frame_ratio:
        shown_ratio = frame_ratio / image_ratio
        crop = (1 - shown_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        shown_ratio = image_ratio / frame_ratio
        crop = (1 - shown_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_header(slide, section, title, subtitle=None):
    add_badge(slide, section.upper(), 0.65, 0.42, max(1.25, min(2.5, len(section) * 0.11 + 0.8)))
    add_text(slide, title, 0.65, 0.82, 12.0, 0.55, 27, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.37, 12.0, 0.42, 13, MUTED)


def add_footer(slide, number):
    add_text(slide, "SD Negeri 232 Maluku Tengah", 0.65, 7.12, 5.5, 0.2, 9, MUTED)
    add_text(slide, str(number), 12.2, 7.12, 0.45, 0.2, 9, MUTED, align=PP_ALIGN.RIGHT)


def new_slide(background=OFF_WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, background)
    return slide


def add_numbered_step(slide, number, title, detail, x, y, width=2.3):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.46), Inches(0.46))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SKY
    circle.line.fill.background()
    frame = circle.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = "Aptos"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    add_text(slide, title, x + 0.58, y - 0.01, width - 0.58, 0.27, 14, NAVY, True)
    add_text(slide, detail, x + 0.58, y + 0.25, width - 0.58, 0.5, 10, MUTED)


# Slide 1 - Cover
slide = new_slide(WHITE)
left_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.25), prs.slide_height)
left_panel.fill.solid()
left_panel.fill.fore_color.rgb = SKY_DARK
left_panel.line.fill.background()
add_picture_contain(slide, LOGO, 0.75, 0.58, 0.75, 0.92)
add_text(slide, "SD NEGERI 232\nMALUKU TENGAH", 1.68, 0.68, 2.75, 0.75, 17, WHITE, True)
add_badge(slide, "PRESENTASI APLIKASI", 0.75, 2.15, 2.15, WHITE, SKY)
add_text(slide, "Sistem Informasi\nJadwal Pelajaran\ndan Nilai Siswa", 0.75, 2.67, 4.0, 2.1, 30, WHITE, True)
add_text(slide, "Pengelolaan akademik sekolah yang lebih tertib, aman, dan mudah digunakan.", 0.75, 5.05, 3.85, 0.85, 16, RGBColor(224, 242, 254))
add_card(slide, 5.8, 0.65, 6.85, 6.15)
add_picture_contain(slide, SCREENSHOTS / "00-login.png", 6.02, 0.88, 6.42, 5.72)
add_text(slide, "Asia/Jayapura - 2026", 0.75, 6.76, 3.0, 0.25, 10, RGBColor(186, 230, 253))

# Slide 2 - Background
slide = new_slide()
add_header(slide, "Latar Belakang", "Mengapa aplikasi ini dibutuhkan?", "Data akademik perlu tersusun agar mudah ditemukan dan tidak saling bertabrakan.")
cards = [
    ("01", "Data tersebar", "Data guru, siswa, kelas, jadwal, dan nilai sulit dipantau jika berada di banyak catatan."),
    ("02", "Risiko kesalahan", "Jadwal dapat bentrok dan nilai dapat tertukar bila proses tidak memiliki validasi."),
    ("03", "Laporan lambat", "Rekap manual membutuhkan waktu dan menyulitkan pemantauan kepala sekolah."),
]
for index, (number, title, detail) in enumerate(cards):
    x = 0.7 + index * 4.2
    add_card(slide, x, 2.15, 3.75, 3.65)
    add_text(slide, number, x + 0.28, 2.42, 0.65, 0.5, 24, SKY, True)
    add_text(slide, title, x + 0.28, 3.12, 3.1, 0.45, 21, NAVY, True)
    add_text(slide, detail, x + 0.28, 3.76, 3.12, 1.35, 15, SLATE)
add_text(slide, "Aplikasi menyatukan proses tersebut dalam satu tempat.", 0.7, 6.28, 11.95, 0.45, 18, SKY_DARK, True, PP_ALIGN.CENTER)
add_footer(slide, 2)

# Slide 3 - Solution
slide = new_slide()
add_header(slide, "Solusi", "Satu aplikasi untuk kegiatan akademik utama", "Pengguna melihat fitur sesuai tanggung jawabnya.")
add_card(slide, 0.7, 2.0, 5.25, 4.38)
add_picture_contain(slide, SCREENSHOTS / "admin-dashboard.png", 0.9, 2.2, 4.85, 3.98)
add_rich_lines(slide, [
    ("Terpusat - ", "data akademik berada dalam satu sistem."),
    ("Terkendali - ", "setiap pengguna memiliki akses yang berbeda."),
    ("Tervalidasi - ", "jadwal bentrok dan nilai tidak sah ditolak."),
    ("Siap laporan - ", "hasil dapat dipantau, dicetak, dan diunduh sebagai PDF."),
], 6.35, 2.1, 5.85, 3.55, 18, 11)
add_text(slide, "Tujuan akhirnya: pekerjaan lebih cepat dan data lebih dapat dipercaya.", 6.35, 5.72, 5.8, 0.65, 18, SKY_DARK, True)
add_footer(slide, 3)

# Slide 4 - Roles
slide = new_slide()
add_header(slide, "Pengguna", "Empat jenis pengguna, empat kebutuhan", "Menu dan data otomatis dibatasi berdasarkan role akun.")
roles = [
    ("ADMIN", "Kelola seluruh data", SKY),
    ("GURU", "Jadwal dan input nilai", GREEN),
    ("SISWA", "Jadwal dan nilai sendiri", AMBER),
    ("KEPALA SEKOLAH", "Monitoring dan laporan", RED),
]
for index, (role, detail, color) in enumerate(roles):
    x = 0.7 + (index % 2) * 6.15
    y = 2.05 + (index // 2) * 2.12
    add_card(slide, x, y, 5.7, 1.62)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.28), Inches(y + 0.32), Inches(0.95), Inches(0.95))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, str(index + 1), x + 0.28, y + 0.46, 0.95, 0.38, 19, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, role, x + 1.5, y + 0.32, 3.8, 0.36, 17, NAVY, True)
    add_text(slide, detail, x + 1.5, y + 0.82, 3.8, 0.42, 14, SLATE)
add_footer(slide, 4)

# Slide 5 - Login
slide = new_slide()
add_header(slide, "Akses", "Masuk dengan akun sekolah", "Registrasi umum dinonaktifkan; akun dibuat dan dikelola oleh Admin.")
add_card(slide, 0.68, 1.9, 8.0, 4.82)
add_picture_contain(slide, SCREENSHOTS / "00-login.png", 0.87, 2.08, 7.62, 4.45)
add_numbered_step(slide, 1, "Isi akun", "Masukkan username atau email.", 9.05, 2.1, 3.45)
add_numbered_step(slide, 2, "Isi password", "Gunakan password yang diberikan sekolah.", 9.05, 3.25, 3.45)
add_numbered_step(slide, 3, "Tekan Masuk", "Dashboard tampil sesuai role pengguna.", 9.05, 4.4, 3.45)
add_text(slide, "Akun nonaktif otomatis ditolak.", 9.05, 5.8, 3.4, 0.45, 15, RED, True)
add_footer(slide, 5)

# Slide 6 - Admin
slide = new_slide()
add_header(slide, "Admin", "Dashboard dan kendali utama", "Admin melihat ringkasan kondisi sekolah dan mengelola seluruh data akademik.")
add_card(slide, 0.7, 1.9, 8.0, 4.82)
add_picture_contain(slide, SCREENSHOTS / "admin-dashboard.png", 0.9, 2.1, 7.6, 4.42)
add_rich_lines(slide, [
    ("Ringkasan - ", "guru, siswa, kelas, mata pelajaran, jadwal, dan nilai."),
    ("Periode aktif - ", "tahun ajaran dan semester yang sedang berjalan."),
    ("Jadwal hari ini - ", "informasi kegiatan akademik terbaru."),
], 9.05, 2.2, 3.35, 2.75, 16, 10)
add_text(slide, "Semua informasi penting terlihat tanpa membuka banyak halaman.", 9.05, 5.25, 3.25, 0.95, 17, SKY_DARK, True)
add_footer(slide, 6)

# Slide 7 - Master data
slide = new_slide()
add_header(slide, "Master Data", "Data dasar yang saling terhubung", "Admin mengisi data sekali, lalu menggunakannya pada jadwal, nilai, dan laporan.")
items = ["Tahun Ajaran", "Semester", "Guru", "Siswa", "Kelas", "Mata Pelajaran", "Jam Pelajaran", "Pengguna"]
for index, item in enumerate(items):
    x = 0.75 + (index % 4) * 3.1
    y = 2.0 + (index // 4) * 1.2
    add_card(slide, x, y, 2.72, 0.88, WHITE)
    add_text(slide, f"{index + 1:02}", x + 0.18, y + 0.23, 0.42, 0.25, 12, SKY, True)
    add_text(slide, item, x + 0.72, y + 0.2, 1.78, 0.35, 14, NAVY, True)
add_card(slide, 0.75, 4.65, 5.92, 1.45, SKY_LIGHT, SKY_LIGHT)
add_text(slide, "Riwayat tetap aman", 1.05, 4.94, 2.5, 0.35, 18, SKY_DARK, True)
add_text(slide, "Data yang sudah dipakai dinonaktifkan, bukan sembarang dihapus.", 1.05, 5.37, 5.2, 0.42, 14, SLATE)
add_card(slide, 6.95, 4.65, 5.65, 1.45, WHITE)
add_text(slide, "Pencarian dan filter", 7.25, 4.94, 2.8, 0.35, 18, NAVY, True)
add_text(slide, "Tabel dapat dicari, diurutkan, disaring, dan dibagi per halaman.", 7.25, 5.37, 4.95, 0.42, 14, SLATE)
add_footer(slide, 7)

# Slide 8 - Schedule
slide = new_slide()
add_header(slide, "Jadwal", "Jadwal mudah dibaca dan bebas bentrok", "Sistem memeriksa konflik sebelum jadwal disimpan.")
add_card(slide, 0.7, 1.9, 7.8, 4.82)
add_picture_contain(slide, SCREENSHOTS / "admin-jadwal-pelajaran.png", 0.9, 2.1, 7.4, 4.42)
rules = [
    ("Bentrok kelas", "Satu kelas tidak boleh memiliki dua pelajaran pada waktu yang sama."),
    ("Bentrok guru", "Satu guru tidak boleh mengajar dua kelas pada waktu yang sama."),
    ("Duplikat", "Jadwal yang sama tidak dapat disimpan dua kali."),
]
for index, (title, detail) in enumerate(rules):
    y = 2.02 + index * 1.45
    add_card(slide, 8.85, y, 3.75, 1.15)
    add_text(slide, str(index + 1), 9.05, y + 0.26, 0.45, 0.28, 14, SKY, True)
    add_text(slide, title, 9.6, y + 0.18, 2.7, 0.3, 15, NAVY, True)
    add_text(slide, detail, 9.6, y + 0.52, 2.65, 0.5, 11, SLATE)
add_footer(slide, 8)

# Slide 9 - Grades
slide = new_slide()
add_header(slide, "Guru", "Input nilai banyak siswa sekaligus", "Guru memilih pengajaran dan bulan, kemudian mengisi Minggu 1 sampai Minggu 4.")
add_card(slide, 0.7, 1.9, 8.35, 4.82)
add_picture_contain(slide, SCREENSHOTS / "guru-input-nilai-terisi.png", 0.9, 2.1, 7.95, 4.42)
add_rich_lines(slide, [
    ("Cepat - ", "seluruh siswa muncul dalam satu tabel."),
    ("Aman - ", "guru hanya dapat mengisi kelas yang diajar."),
    ("Valid - ", "nilai harus 0-100; kosong berarti belum dinilai."),
    ("Otomatis - ", "rata-rata menghitung nilai yang sudah terisi saja."),
], 9.35, 2.15, 3.1, 3.65, 15, 9)
add_footer(slide, 9)

# Slide 10 - Student
slide = new_slide()
add_header(slide, "Siswa", "Siswa melihat jadwal dan nilainya sendiri", "Akses siswa dibatasi otomatis berdasarkan akun dan kelas aktif.")
add_card(slide, 0.72, 1.95, 5.95, 4.72)
add_picture_contain(slide, SCREENSHOTS / "siswa-jadwal-pelajaran.png", 0.92, 2.15, 5.55, 4.32)
add_card(slide, 6.92, 1.95, 5.7, 4.72)
add_picture_contain(slide, SCREENSHOTS / "siswa-nilai-saya.png", 7.12, 2.15, 5.3, 4.32)
add_badge(slide, "JADWAL KELAS", 1.05, 6.15, 1.55)
add_badge(slide, "NILAI SAYA", 7.25, 6.15, 1.4)
add_footer(slide, 10)

# Slide 11 - Headmaster
slide = new_slide()
add_header(slide, "Kepala Sekolah", "Monitoring tanpa mengubah data", "Kepala Sekolah memperoleh informasi untuk pengawasan dan pengambilan keputusan.")
add_card(slide, 0.7, 1.9, 7.95, 4.82)
add_picture_contain(slide, SCREENSHOTS / "kepala-dashboard.png", 0.9, 2.1, 7.55, 4.42)
add_rich_lines(slide, [
    ("Pantau - ", "jumlah data dan aktivitas akademik."),
    ("Filter - ", "jadwal dan nilai berdasarkan kelas, guru, siswa, atau periode."),
    ("Laporan - ", "cetak dan unduh PDF saat diperlukan."),
    ("Read-only - ", "tidak memiliki CRUD penuh seperti Admin."),
], 9.0, 2.18, 3.25, 3.7, 16, 10)
add_footer(slide, 11)

# Slide 12 - Reports
slide = new_slide()
add_header(slide, "Laporan", "Preview, cetak, dan unduh PDF", "Informasi sekolah dan periode ditampilkan pada laporan resmi.")
add_card(slide, 0.7, 1.9, 6.0, 4.82)
add_picture_contain(slide, SCREENSHOTS / "admin-laporan-jadwal.png", 0.9, 2.1, 5.6, 4.42)
add_card(slide, 6.95, 1.9, 5.68, 4.82)
add_picture_contain(slide, SCREENSHOTS / "admin-laporan-nilai.png", 7.15, 2.1, 5.28, 4.42)
add_badge(slide, "LAPORAN JADWAL", 0.98, 6.15, 1.8)
add_badge(slide, "LAPORAN NILAI", 7.22, 6.15, 1.7)
add_footer(slide, 12)

# Slide 13 - Mobile
slide = new_slide()
add_header(slide, "Responsif", "Nyaman digunakan dari telepon genggam", "Menu berubah menjadi drawer dan tabel tetap dapat digeser secara horizontal.")
phones = [
    ("Guru", SCREENSHOTS / "mobile-guru-dashboard.png"),
    ("Input Nilai", SCREENSHOTS / "mobile-guru-input-nilai-terisi.png"),
    ("Siswa", SCREENSHOTS / "mobile-siswa-nilai-saya.png"),
]
for index, (label, path) in enumerate(phones):
    x = 1.0 + index * 4.05
    add_card(slide, x, 1.9, 3.3, 4.9, WHITE)
    add_picture_contain(slide, path, x + 0.17, 2.08, 2.96, 4.32)
    add_badge(slide, label.upper(), x + 0.82, 6.35, 1.65)
add_footer(slide, 13)

# Slide 14 - Security
slide = new_slide()
add_header(slide, "Keamanan", "Akses dan integritas data dijaga", "Perlindungan diterapkan di backend, bukan hanya menyembunyikan menu.")
security = [
    ("Role & permission", "Admin, Guru, Siswa, dan Kepala Sekolah memiliki izin berbeda."),
    ("Pencegahan IDOR", "Guru dan siswa hanya mengambil data yang terkait dengan akun login."),
    ("Validasi server", "Jadwal, nilai, akun, dan upload diperiksa sebelum disimpan."),
    ("Password aman", "Password di-hash dan tidak pernah ditampilkan sebagai teks biasa."),
    ("Riwayat akademik", "Foreign key dan status nonaktif menjaga data historis."),
    ("Pengujian otomatis", "18 test dan 76 assertion memastikan aturan penting berjalan."),
]
for index, (title, detail) in enumerate(security):
    x = 0.7 + (index % 2) * 6.15
    y = 1.95 + (index // 2) * 1.55
    add_card(slide, x, y, 5.7, 1.18)
    add_text(slide, title, x + 0.28, y + 0.2, 2.25, 0.32, 15, NAVY, True)
    add_text(slide, detail, x + 2.45, y + 0.18, 2.9, 0.72, 11, SLATE)
add_footer(slide, 14)

# Slide 15 - Demo flow
slide = new_slide()
add_header(slide, "Demo", "Urutan demonstrasi yang disarankan", "Alur ini dapat disampaikan dalam sekitar 7-10 menit.")
steps = [
    ("Login", "Masuk sebagai Admin."),
    ("Dashboard", "Jelaskan ringkasan sekolah."),
    ("Jadwal", "Tunjukkan filter dan validasi bentrok."),
    ("Nilai", "Masuk sebagai Guru dan isi nilai massal."),
    ("Hasil", "Masuk sebagai Siswa, lalu buka laporan PDF."),
]
for index, (title, detail) in enumerate(steps):
    x = 0.75 + index * 2.48
    add_card(slide, x, 2.35, 2.15, 2.65)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.68), Inches(2.68), Inches(0.78), Inches(0.78))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SKY
    circle.line.fill.background()
    add_text(slide, str(index + 1), x + 0.68, 2.87, 0.78, 0.28, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.18, 3.7, 1.8, 0.35, 16, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, detail, x + 0.18, 4.18, 1.8, 0.58, 11, SLATE, align=PP_ALIGN.CENTER)
    if index < len(steps) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.17), Inches(3.23), Inches(0.3), Inches(0.35))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(186, 230, 253)
        line.line.fill.background()
add_card(slide, 1.4, 5.55, 10.5, 0.72, SKY_LIGHT, SKY_LIGHT)
add_text(slide, "Fokus pada manfaat dan alur kerja. Hindari penjelasan teknis yang terlalu dalam.", 1.65, 5.78, 10.0, 0.3, 15, SKY_DARK, True, PP_ALIGN.CENTER)
add_footer(slide, 15)

# Slide 16 - Closing
slide = new_slide(SKY_DARK)
add_picture_contain(slide, LOGO, 0.75, 0.72, 1.0, 1.22)
add_text(slide, "SD NEGERI 232 MALUKU TENGAH", 2.0, 0.98, 5.8, 0.42, 17, WHITE, True)
add_text(slide, "Jadwal lebih tertib.\nNilai lebih mudah.\nLaporan lebih cepat.", 0.78, 2.35, 7.5, 2.15, 34, WHITE, True)
add_text(slide, "Sistem Informasi Jadwal Pelajaran dan Nilai Siswa", 0.78, 4.92, 6.7, 0.6, 18, RGBColor(224, 242, 254))
add_card(slide, 8.55, 1.5, 3.9, 4.3, WHITE, WHITE)
add_text(slide, "Terima Kasih", 8.95, 2.18, 3.1, 0.55, 27, NAVY, True, PP_ALIGN.CENTER)
add_text(slide, "Pertanyaan dan diskusi", 8.95, 2.95, 3.1, 0.42, 16, SKY_DARK, True, PP_ALIGN.CENTER)
add_text(slide, "Aplikasi siap digunakan untuk mendukung pengelolaan akademik sekolah.", 9.1, 3.75, 2.8, 1.15, 15, SLATE, align=PP_ALIGN.CENTER)
add_text(slide, "SISD 232", 9.45, 5.05, 2.1, 0.38, 14, MUTED, True, PP_ALIGN.CENTER)

prs.core_properties.title = "Presentasi Aplikasi SISD 232"
prs.core_properties.subject = "Sistem Informasi Jadwal Pelajaran dan Nilai Siswa"
prs.core_properties.author = "SD Negeri 232 Maluku Tengah"
prs.core_properties.keywords = "sekolah, jadwal, nilai, Laravel, SISD 232"
prs.save(OUTPUT)

print(f"Presentasi dibuat: {OUTPUT}")
print(f"Jumlah slide: {len(prs.slides)}")
